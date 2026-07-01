"""PowerPoint (.pptx) → Markdown converter using python-pptx.

Each slide becomes a section with text, tables, and images.
"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from converters.base import BaseConverter, ConversionResult


class PPTXConverter(BaseConverter):
    """Convert PowerPoint presentations to Markdown with image extraction."""

    def convert(self) -> ConversionResult:
        prs = Presentation(str(self.file_path))

        md_parts: list[str] = []
        images: list[tuple[str, bytes]] = []

        # Presentation title
        title = self.file_path.stem
        if prs.core_properties.title:
            title = prs.core_properties.title
        md_parts.append(f"# {title}\n")

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_md, slide_images = self._convert_slide(slide, slide_num)
            if slide_md.strip():
                md_parts.append(slide_md)
            images.extend(slide_images)

        metadata = {
            "format": "pptx",
            "slide_count": len(prs.slides),
            "title": title,
            "author": prs.core_properties.author or "",
        }

        return ConversionResult(markdown="\n\n".join(md_parts), images=images, metadata=metadata)

    def _convert_slide(
        self, slide, slide_num: int
    ) -> tuple[str, list[tuple[str, bytes]]]:
        """Convert one slide to markdown."""
        lines: list[str] = []
        images: list[tuple[str, bytes]] = []

        lines.append(f"---")
        lines.append(f"## Slide {slide_num}")
        lines.append("")

        # Slide title from layout placeholder
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        if title_text:
            lines.append(f"### {title_text}")
            lines.append("")

        # Process shapes
        for shape in slide.shapes:
            shape_lines, shape_images = self._process_shape(shape)
            lines.extend(shape_lines)
            images.extend(shape_images)

        # Speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append("> **备注:** " + notes.replace("\n", "\n> "))

        return "\n".join(lines), images

    def _process_shape(self, shape) -> tuple[list[str], list[tuple[str, bytes]]]:
        """Dispatch shape to the right handler based on type."""
        lines: list[str] = []
        images: list[tuple[str, bytes]] = []

        # Group shapes — recurse
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                cl, ci = self._process_shape(child)
                lines.extend(cl)
                images.extend(ci)
            return lines, images

        # Pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_bytes = shape.image.blob
                ext = shape.image.content_type.split("/")[-1]
                filename, png_bytes = self._next_image(img_bytes, ext)
                images.append((filename, png_bytes))
                lines.append(self._image_md(filename, alt=shape.name))
                lines.append("")
            except Exception:
                pass
            return lines, images

        # Tables
        if shape.has_table:
            tbl_lines = self._table_to_md(shape.table)
            lines.extend(tbl_lines)
            lines.append("")
            return lines, images

        # Text frames
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)
                lines.append("")

        return lines, images

    def _table_to_md(self, table) -> list[str]:
        """Convert a PPTX table to pipe-table markdown."""
        lines: list[str] = []
        max_cols = len(table.columns)
        if max_cols == 0:
            return []

        is_first = True
        for row in table.rows:
            cells = [self._escape_pipe(cell.text) for cell in row.cells]
            line = "| " + " | ".join(cells) + " |"
            lines.append(line)
            if is_first:
                lines.append("| " + " | ".join("---" for _ in range(max_cols)) + " |")
                is_first = False

        return lines
